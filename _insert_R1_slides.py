"""Insert two R1-result slides into 'Sakharov Midway v2.pptx' right after the
existing 'Results' tabs (slide 37). Two new slides:
  Slide A: Frame-level P/R/F1 @ IoU=0.3 at best-F1 threshold
  Slide B: Track-level metrics at the same threshold
"""
import json
from pathlib import Path
from copy import deepcopy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = Path("/home/dsa/stenosis/Sakharov Midway v2.pptx")
METRICS_PATH = Path(
    "/home/dsa/stenosis/rfdetr_video/runs/video_overfit_R1/intuitive_metrics.json"
)

INSERT_AFTER_SLIDE = 37  # 1-indexed; new slides become 38, 39

# ---- Load metrics ----
with open(METRICS_PATH) as f:
    M = json.load(f)
reps = {r["dataset"]: r for r in M["reports"]}
ds2 = reps["dataset2_split_test"]
cad = reps["cadica_50plus_new"]


def pct(x, total):
    return f"{x}  ({x / total:.1%})"


# ---- Build slide content ----
# Both columns use the best-F1 threshold per dataset (different per dataset).

ds2_b = ds2["best_F1_at_IoU0.3"]
cad_b = cad["best_F1_at_IoU0.3"]
ds2_t = ds2["tracks_at_best_thr"]
cad_t = cad["tracks_at_best_thr"]

FRAME_HEADER = [
    "Metric (IoU = 0.3)",
    f"dataset2_split_test  (score ≥ {ds2_b['thr']:.2f})",
    f"cadica_50plus_new  (score ≥ {cad_b['thr']:.2f})",
]
FRAME_ROWS = [
    ["Precision",      f"{ds2_b['P']:.3f}",  f"{cad_b['P']:.3f}"],
    ["Recall",         f"{ds2_b['R']:.3f}",  f"{cad_b['R']:.3f}"],
    ["F1",             f"{ds2_b['F1']:.3f}", f"{cad_b['F1']:.3f}"],
    ["AP @ 0.3",       f"{ds2['AP@0.3']:.3f}", f"{cad['AP@0.3']:.3f}"],
    ["AP @ 0.5",       f"{ds2['AP@0.5']:.3f}", f"{cad['AP@0.5']:.3f}"],
    ["TP / FP / FN",   f"{ds2_b['TP']} / {ds2_b['FP']} / {ds2_b['FN']}",
                       f"{cad_b['TP']} / {cad_b['FP']} / {cad_b['FN']}"],
    ["# videos / frames / GT boxes",
        f"{ds2['n_videos']} / {ds2['n_frames']} / {ds2['n_gt_total']}",
        f"{cad['n_videos']} / {cad['n_frames']} / {cad['n_gt_total']}"],
]

TRACK_HEADER = [
    "Track-level metric (match IoU = 0.3)",
    f"dataset2_split_test  (score ≥ {ds2_b['thr']:.2f})",
    f"cadica_50plus_new  (score ≥ {cad_b['thr']:.2f})",
]
TRACK_ROWS = [
    ["Total GT tracks (stenosis instances)",
        f"{ds2_t['n_tracks']}", f"{cad_t['n_tracks']}"],
    ["Ever detected  (≥1 hit)",
        pct(ds2_t['tracks_ever_found'],     ds2_t['n_tracks']),
        pct(cad_t['tracks_ever_found'],     cad_t['n_tracks'])],
    ["Well-detected  (≥50% frames)",
        pct(ds2_t['tracks_well_found'],     ds2_t['n_tracks']),
        pct(cad_t['tracks_well_found'],     cad_t['n_tracks'])],
    ["Excellently-detected  (≥80% frames)",
        pct(ds2_t['tracks_excellently_found'], ds2_t['n_tracks']),
        pct(cad_t['tracks_excellently_found'], cad_t['n_tracks'])],
    ["Completely missed",
        pct(ds2_t['tracks_missed'],         ds2_t['n_tracks']),
        pct(cad_t['tracks_missed'],         cad_t['n_tracks'])],
    ["Mean per-track recall",
        f"{ds2_t['mean_per_track_recall']:.3f}",
        f"{cad_t['mean_per_track_recall']:.3f}"],
    ["Mean longest-streak ratio per track",
        f"{ds2_t['mean_longest_streak_ratio']:.3f}",
        f"{cad_t['mean_longest_streak_ratio']:.3f}"],
    ["Mean fragmentations per track",
        f"{ds2_t['mean_frags_per_track']:.3f}",
        f"{cad_t['mean_frags_per_track']:.3f}"],
]


# ---- Build slides ----

prs = Presentation(str(PPTX_PATH))
layout = prs.slide_layouts[2]  # TITLE_AND_TWO_COLUMNS, same as existing Results
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# Reference geometry from existing Results slide 33:
#   title placeholder: left=0.787in top=0.118in w=8.425 h=0.648
#   table:             left=0.787in top=0.900in w=8.075 h=4.667
TITLE_LEFT, TITLE_TOP = Inches(0.787), Inches(0.118)
TITLE_W, TITLE_H = Inches(8.425), Inches(0.648)
SUB_TOP, SUB_H = Inches(0.78), Inches(0.45)
TBL_LEFT, TBL_TOP = Inches(0.787), Inches(1.30)
TBL_W, TBL_H = Inches(8.075), Inches(4.25)


def _set_title(slide, text):
    """Set the title placeholder text."""
    title_ph = None
    if slide.shapes.title is not None:
        title_ph = slide.shapes.title
    else:
        for shape in slide.placeholders:
            if "Title" in (shape.placeholder_format.idx, shape.name):
                title_ph = shape
                break
    if title_ph is None:
        # Create a textbox if no title placeholder
        title_ph = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = title_ph.text_frame
    tf.text = text
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(28)
            run.font.bold = True


def _add_subtitle(slide, text):
    box = slide.shapes.add_textbox(TITLE_LEFT, SUB_TOP, TITLE_W, SUB_H)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def _add_table(slide, header, rows, highlight_rows=None):
    """Add a table whose first row is header, rest are body. highlight_rows is
    list of body-row indices (0-based, not counting header) to bold."""
    highlight_rows = set(highlight_rows or [])
    n_rows = 1 + len(rows)
    n_cols = len(header)
    table_shape = slide.shapes.add_table(
        n_rows, n_cols, TBL_LEFT, TBL_TOP, TBL_W, TBL_H,
    )
    tbl = table_shape.table
    # Set first col wider, other two equal
    col_widths_in = [3.0, 2.55, 2.55]
    total_emu = sum(int(w * 914400) for w in col_widths_in)
    scale = TBL_W / total_emu
    for c, w in enumerate(col_widths_in):
        tbl.columns[c].width = Emu(int(w * 914400 * scale))

    def _fill_cell(cell, text, *, bold=False, header=False, size=14, align=None):
        tf = cell.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        if align is not None:
            p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold or header
        if header:
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        else:
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # Header
    for c, h in enumerate(header):
        _fill_cell(tbl.cell(0, c), h, header=True, size=13,
                   align=PP_ALIGN.CENTER)
    # Body
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            is_hi = r in highlight_rows
            align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            _fill_cell(tbl.cell(r + 1, c), val, bold=is_hi, align=align,
                       size=14 if c == 0 else 15)

    # Zebra striping on body rows
    for r in range(1, n_rows):
        if (r - 1) % 2 == 1:
            for c in range(n_cols):
                cell = tbl.cell(r, c)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)


# Slide A: frame-level
slide_a = prs.slides.add_slide(layout)
_set_title(slide_a, "Results")
_add_subtitle(
    slide_a,
    "video_overfit_R1  —  frame-level Precision / Recall / F1 @ IoU = 0.3, "
    "best-F1 score threshold (per dataset)",
)
_add_table(slide_a, FRAME_HEADER, FRAME_ROWS, highlight_rows=[2])  # F1 row bold

# Slide B: track-level
slide_b = prs.slides.add_slide(layout)
_set_title(slide_b, "Results")
_add_subtitle(
    slide_b,
    "video_overfit_R1  —  track-level coverage & fragmentation, "
    "match IoU = 0.3, same best-F1 thresholds",
)
_add_table(slide_b, TRACK_HEADER, TRACK_ROWS, highlight_rows=[1, 7])  # ever + frag


# ---- Reorder: place these two right after slide 37 ----
sld_lst = prs.slides._sldIdLst
sld_elements = list(sld_lst)
# Newly added slides are the last two
new_a, new_b = sld_elements[-2], sld_elements[-1]
sld_lst.remove(new_a)
sld_lst.remove(new_b)
# After removal, insertion index = INSERT_AFTER_SLIDE in 1-indexed = same 0-indexed
sld_lst.insert(INSERT_AFTER_SLIDE,     new_a)
sld_lst.insert(INSERT_AFTER_SLIDE + 1, new_b)


prs.save(str(PPTX_PATH))
print(f"Saved → {PPTX_PATH}")
print(f"  Inserted 2 new slides after original slide {INSERT_AFTER_SLIDE}")
print(f"  Total slides: {len(prs.slides)}")
