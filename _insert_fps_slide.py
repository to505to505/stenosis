"""Insert an 'Inference speed' slide right after the track-level R1 slide
(file position 40, between slide 39 stickiness and slide 40 Conclusion)."""
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX = Path("/home/dsa/stenosis/Sakharov Midway v2.pptx")
BENCH = Path("/home/dsa/stenosis/rfdetr_video/runs/video_overfit_R1/speed_benchmark.json")

bench = json.loads(BENCH.read_text())
amp_fps = bench["amp"]["fps"]
amp_lat = bench["amp"]["latency_ms_mean"]
amp_p95 = bench["amp"]["latency_ms_p95"]
fp32_fps = bench["fp32"]["fps"]
fp32_lat = bench["fp32"]["latency_ms_mean"]
gpu = bench["gpu"]
T = bench["T"]
size = bench["img_size"]
params = bench["params_M"]

prs = Presentation(str(PPTX))
# Same layout as the other Results slides
layout = prs.slide_layouts[2]
slide = prs.slides.add_slide(layout)

# ── Title placeholder ────────────────────────────────────────────────────────
# The TITLE_AND_TWO_COLUMNS layout auto-fills the title placeholder with empty
# text; set it directly via the placeholder collection.
title_ph = None
for ph in slide.placeholders:
    try:
        if ph.placeholder_format.type == 1:  # PP_PLACEHOLDER.TITLE
            title_ph = ph
            break
    except Exception:
        pass
if title_ph is None:
    title_ph = slide.shapes.title
if title_ph is not None:
    title_ph.text_frame.text = "Results"
    for p in title_ph.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True

# ── Subtitle ─────────────────────────────────────────────────────────────────
TITLE_LEFT = Inches(0.787)
TITLE_W = Inches(8.425)

box = slide.shapes.add_textbox(TITLE_LEFT, Inches(0.78), TITLE_W, Inches(0.50))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = (
    f"video_overfit_R1 — inference speed on {gpu}, "
    f"T = {T} × {size}² window, AMP fp16, batch = 1"
)
run.font.size = Pt(13)
run.font.italic = True
run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ── Big headline FPS ─────────────────────────────────────────────────────────
big_box = slide.shapes.add_textbox(
    Inches(0.787), Inches(1.55), Inches(8.425), Inches(1.40),
)
btf = big_box.text_frame
btf.word_wrap = True
bp = btf.paragraphs[0]
bp.alignment = PP_ALIGN.CENTER
br = bp.add_run()
br.text = f"{amp_fps:.1f} FPS"
br.font.size = Pt(72)
br.font.bold = True
br.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

# Latency subtext
bp2 = btf.add_paragraph()
bp2.alignment = PP_ALIGN.CENTER
br2 = bp2.add_run()
br2.text = f"{amp_lat:.1f} ms  per centre-frame  (p95 {amp_p95:.1f} ms)"
br2.font.size = Pt(18)
br2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ── Comparison table ─────────────────────────────────────────────────────────
TBL_LEFT = Inches(1.40)
TBL_TOP = Inches(3.30)
TBL_W = Inches(7.20)
TBL_H = Inches(1.10)

rows = [
    ["Configuration", "Per-frame latency", "FPS"],
    ["AMP fp16 (training-style, on RTX 3060)",
        f"{amp_lat:.1f} ms",     f"{amp_fps:.1f}"],
    ["FP32  (no AMP, same hardware)",
        f"{fp32_lat:.1f} ms",    f"{fp32_fps:.1f}"],
]
table_shape = slide.shapes.add_table(3, 3, TBL_LEFT, TBL_TOP, TBL_W, TBL_H)
tbl = table_shape.table
col_widths_in = [4.00, 1.80, 1.40]
total = sum(int(w * 914400) for w in col_widths_in)
scale = TBL_W / total
for c, w in enumerate(col_widths_in):
    tbl.columns[c].width = Emu(int(w * 914400 * scale))


def _fill(cell, text, *, header=False, bold=False, size=14, align=None):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold or header
    if header:
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    else:
        r.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


for c, h in enumerate(rows[0]):
    _fill(tbl.cell(0, c), h, header=True, size=13, align=PP_ALIGN.CENTER)
for r, row in enumerate(rows[1:], start=1):
    for c, val in enumerate(row):
        align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        # Bold the AMP row (row index 1, the headline configuration)
        bold = (r == 1)
        _fill(tbl.cell(r, c), val, bold=bold, align=align,
              size=14 if c == 0 else 15)
# Zebra: row 2 stripe
for c in range(3):
    cell = tbl.cell(2, c)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)

# ── Footer note ──────────────────────────────────────────────────────────────
foot = slide.shapes.add_textbox(
    Inches(0.787), Inches(4.65), Inches(8.425), Inches(0.80),
)
ftf = foot.text_frame
ftf.word_wrap = True
fp = ftf.paragraphs[0]
fp.alignment = PP_ALIGN.LEFT
fr = fp.add_run()
fr.text = (
    f"Clinical coronary angiography is acquired at ~15 FPS  →  "
    f"~{amp_fps/15:.1f}× real-time headroom on a single consumer GPU.  "
    f"Model: {params:.0f} M parameters."
)
fr.font.size = Pt(13)
fr.font.italic = True
fr.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ── Reorder: move the new slide to file position 40 ──────────────────────────
sld_lst = prs.slides._sldIdLst
els = list(sld_lst)
new_sld = els[-1]
sld_lst.remove(new_sld)
sld_lst.insert(39, new_sld)   # 0-indexed 39 → 1-indexed position 40

prs.save(str(PPTX))

# Verify
prs = Presentation(str(PPTX))
print(f"Total slides: {len(prs.slides)}")
print("\nSlide 40 shapes:")
s = prs.slides[39]
for sh in s.shapes:
    info = f"  {sh.shape_type} ph={sh.is_placeholder}"
    if sh.has_text_frame and sh.text_frame.text.strip():
        info += f" '{sh.text_frame.text[:100]}'"
    if sh.has_table:
        info += f" TABLE {len(sh.table.rows)}x{len(sh.table.columns)}"
    print(info)
