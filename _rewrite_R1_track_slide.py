"""Rewrite the track-level R1 slide (slide 39) to lead with the
'stickiness' story (once a stenosis is detected, the model keeps it
locked) and de-emphasize the lower 'ever-detected' coverage.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX = Path("/home/dsa/stenosis/Sakharov Midway v2.pptx")
SLIDE_NUM = 39  # 1-indexed

prs = Presentation(str(PPTX))
slide = prs.slides[SLIDE_NUM - 1]

# ── Strip everything except the title placeholder ────────────────────────────
title_shape = slide.shapes.title
to_remove = [s for s in slide.shapes if s is not title_shape]
for s in to_remove:
    s._element.getparent().remove(s._element)

# ── Subtitle ─────────────────────────────────────────────────────────────────
TITLE_LEFT = Inches(0.787)
SUB_TOP = Inches(0.78)
TITLE_W = Inches(8.425)
box = slide.shapes.add_textbox(TITLE_LEFT, SUB_TOP, TITLE_W, Inches(0.50))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = (
    "video_overfit_R1  —  temporal tracking stability: "
    "once a stenosis is detected, the model keeps it locked"
)
run.font.size = Pt(13)
run.font.italic = True
run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ── Table ────────────────────────────────────────────────────────────────────
TBL_LEFT = Inches(0.787)
TBL_TOP = Inches(1.35)
TBL_W = Inches(8.075)
TBL_H = Inches(4.10)

HEADER = [
    "Track-level metric  (match IoU = 0.3)",
    "dataset2_split_test  (score ≥ 0.52)",
    "cadica_50plus_new  (score ≥ 0.47)",
]
# Lead with stickiness numbers; coverage row at the bottom, dimmed.
ROWS = [
    ["Mean per-track recall  (among detected)",                 "0.83",                 "0.84"],
    ["Mean longest correct streak / track length  (among detected)", "0.80",            "0.81"],
    ["Tracks held ≥50% of frames  (% of detected)",             "87.5%   (91 / 104)",   "88.1%   (288 / 327)"],
    ["Tracks held ≥80% of frames  (% of detected)",             "71.2%   (74 / 104)",   "71.6%   (234 / 327)"],
    ["Mean fragmentations per detected track",                  "0.43",                 "0.16"],
    ["Detection coverage  (tracks with ≥1 hit)",                "104 / 176",            "327 / 795"],
]
HIGHLIGHT_BODY_ROWS = {0, 1, 4}     # bold the 3 strongest stickiness rows
DIM_BODY_ROWS = {5}                  # last row: coverage, dimmed

n_rows = 1 + len(ROWS)
n_cols = len(HEADER)
table_shape = slide.shapes.add_table(n_rows, n_cols, TBL_LEFT, TBL_TOP, TBL_W, TBL_H)
tbl = table_shape.table

# Column widths
col_widths_in = [3.4, 2.35, 2.35]
total = sum(int(w * 914400) for w in col_widths_in)
scale = TBL_W / total
for c, w in enumerate(col_widths_in):
    tbl.columns[c].width = Emu(int(w * 914400 * scale))


def _fill(cell, text, *, header=False, bold=False, dim=False,
          size=14, align=None):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold or header
    if header:
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
    elif dim:
        run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        run.font.italic = True
    else:
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# Header
for c, h in enumerate(HEADER):
    _fill(tbl.cell(0, c), h, header=True, size=13, align=PP_ALIGN.CENTER)

# Body
for r, row in enumerate(ROWS):
    is_hi = r in HIGHLIGHT_BODY_ROWS
    is_dim = r in DIM_BODY_ROWS
    for c, val in enumerate(row):
        align = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        _fill(tbl.cell(r + 1, c), val,
              bold=is_hi, dim=is_dim,
              align=align,
              size=14 if c == 0 else (15 if is_hi else 14))

# Zebra stripes (skip dimmed row)
for r in range(1, n_rows):
    body_idx = r - 1
    if body_idx in DIM_BODY_ROWS:
        continue
    if body_idx % 2 == 1:
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF6, 0xFA)

# ── Caption under the table ──────────────────────────────────────────────────
cap = slide.shapes.add_textbox(
    TITLE_LEFT, Inches(5.20), Inches(8.425), Inches(0.40),
)
ctf = cap.text_frame
ctf.word_wrap = True
cp = ctf.paragraphs[0]
cp.alignment = PP_ALIGN.LEFT
crun = cp.add_run()
crun.text = (
    "Among detected stenoses the model holds ~83–84% of frames "
    "on average, with a single uninterrupted streak covering "
    "~80% of the track length and 0.16–0.43 break events per track."
)
crun.font.size = Pt(11)
crun.font.italic = True
crun.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

prs.save(str(PPTX))
print(f"Rewrote slide {SLIDE_NUM} in {PPTX}")
