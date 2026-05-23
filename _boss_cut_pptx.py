"""Boss-version of 'Sakharov Midway v2.pptx':

  - Hide thesis-y slides (research questions, basic medical "what is stenosis",
    rationale "Explanations").
  - Hide duplicate Results placeholder slides (34-37; the first Results table
    on 33 stays).
  - Hide the segmentation-chapter appendix (originals 41-88).
  - Hide template-junk slides after Thank You (originals 90-128).
  - Move 'Thank you' slide right after the (new) Conclusion slide.

A hidden slide has ``show="0"`` on its ``<p:sld>`` element; it is skipped in
slideshow mode but kept in the file for reference.
"""
from pptx import Presentation
from pathlib import Path

PPTX_PATH = Path("/home/dsa/stenosis/Sakharov Midway v2.pptx")

prs = Presentation(str(PPTX_PATH))
assert len(prs.slides) == 128, f"Expected 128, got {len(prs.slides)}"

# ── Step 1: Move 'Thank you' (slide 89) to position 41 (right after slide 40
#            Conclusion) for clean edit view.
sld_lst = prs.slides._sldIdLst
els = list(sld_lst)
thank_you_idx = 88           # 0-indexed → slide 89
thank_you = els[thank_you_idx]
# Verify it's actually the Thank You slide
ty_text = prs.slides[thank_you_idx].shapes.title.text \
    if prs.slides[thank_you_idx].shapes.title is not None else ""
assert "Thank" in ty_text, f"Slide 89 isn't Thank You, got: {ty_text!r}"
sld_lst.remove(thank_you)
sld_lst.insert(40, thank_you)   # 0-indexed 40 → position 41

# After move, current positions (1-indexed):
#   1..40            → originals 1..40
#   41               → Thank you
#   42..89           → originals 41..88   (shifted +1)
#   90..128          → originals 90..128

# ── Step 2: Hide the unwanted slides (1-indexed positions in the post-move
#            order).
HIDE = (
    [2, 3, 4]                      # "what is stenosis" medical intro
    + [5, 8]                       # Primary RQ & Research Questions
    + [9, 10, 11, 12, 13]          # "Explanations" — thesis design rationale
    + [34, 35, 36, 37]             # duplicate Results placeholders
    + list(range(42, 90))          # appendix + segmentation chapter
    + list(range(90, 129))         # template-garbage appendix
)
# 89 is now empty after the move — re-check the final list of positions:
# positions 42-89 were originals 41-88
# position 90+ are originals 90+
HIDE = sorted(set(HIDE))

for pos in HIDE:
    slide = prs.slides[pos - 1]
    slide._element.set("show", "0")

# Visible slides (sanity)
visible = [i + 1 for i, s in enumerate(prs.slides)
           if s._element.get("show") != "0"]
hidden = [i + 1 for i, s in enumerate(prs.slides)
          if s._element.get("show") == "0"]

prs.save(str(PPTX_PATH))

print(f"Total slides:   {len(prs.slides)}")
print(f"Visible:        {len(visible)}  → {visible}")
print(f"Hidden:         {len(hidden)}")
print()
print("Visible slide titles:")
for pos in visible:
    s = prs.slides[pos - 1]
    title = ""
    if s.shapes.title is not None:
        title = s.shapes.title.text.replace("\n", " ").strip()[:90]
    else:
        for shape in s.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                title = shape.text_frame.text.replace("\n", " ").strip()[:90]
                break
    print(f"  {pos:3d}: {title}")
